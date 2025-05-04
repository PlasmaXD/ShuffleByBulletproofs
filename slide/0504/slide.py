import os
from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "CMS + LDP 手法の概要"
subtitle.text = "Count–Min Sketch とローカル差分プライバシーの組み合わせ"

# Slide 2: Count–Min Sketch (CMS) 概要
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.shapes.placeholders[1]
title.text = "Count–Min Sketch (CMS) の概要"
tf = body.text_frame
tf.text = "• d×w 多重ハッシュ行列\n• 要素追加: 各行で対応バケットをインクリメント\n• 頻度推定: 各行の値の最小値を返す"
for paragraph in tf.paragraphs:
    paragraph.font.size = Pt(18)

# Slide 3: ローカル差分プライバシー (LDP) 基礎
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.shapes.placeholders[1]
title.text = "ローカル差分プライバシー (LDP) の基礎"
tf = body.text_frame
tf.text = "• 各ユーザーが送信前にデータをランダム化\n• ε-LDP: 出力確率の比が e^ε 以下\n• Randomized Response が代表的手法"
for paragraph in tf.paragraphs:
    paragraph.font.size = Pt(18)

# Slide 4: クライアント側のランダム化手順
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.shapes.placeholders[1]
title.text = "クライアント側のランダム化手順"
tf = body.text_frame
tf.text = "1. アイテム x を d 本のハッシュ関数でマッピング\n2. 元のビットに対して瞬間ランダム応答(IRR)を適用\n3. ノイズ付きビットをサーバーに送信"
for paragraph in tf.paragraphs:
    paragraph.font.size = Pt(18)

# Slide 5: サーバー側の集約と推定
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.shapes.placeholders[1]
title.text = "サーバー側の集約と推定"
tf = body.text_frame
tf.text = "• 受信ビットをSketchセルに加算\n• 観測率から逆ノイズ化: (y-p)/(q-p)\n• 各行の最小値を最終推定値とする"
for paragraph in tf.paragraphs:
    paragraph.font.size = Pt(18)

# Slide 6: 理論保証と誤差特性
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.shapes.placeholders[1]
title.text = "理論保証と誤差特性"
tf = body.text_frame
tf.text = "• ε-LDP 合成則によりプライバシー保証\n• CMSのオーバーエスティメーション＋LDPノイズ\n• 適切な w, d 選択で O((1/ε)√log(1/δ)) 誤差"
for paragraph in tf.paragraphs:
    paragraph.font.size = Pt(18)

# Slide 7: 応用例
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.shapes.placeholders[1]
title.text = "応用例"
tf = body.text_frame
tf.text = "• Apple: 絵文字・語彙統計収集\n• 大規模ログ解析と範囲クエリ\n• ストリームデータの頻度推定"
for paragraph in tf.paragraphs:
    paragraph.font.size = Pt(18)

# Slide 8: 参考文献
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.shapes.placeholders[1]
title.text = "参考文献"
tf = body.text_frame
references = [
    "Cormode et al., VLDB 2020",
    "Apple, Learning with Privacy at Scale",
    "NeurIPS 2022: Differentially Private Linear Sketches",
    "Wikipedia: Local Differential Privacy",
    "GeeksforGeeks: Count–Min Sketch"
]
for ref in references:
    p = tf.add_paragraph()
    p.text = f"• {ref}"
    p.level = 0
    p.font.size = Pt(14)
tf.paragraphs[0].text = ""  # Remove initial empty bullet

# Save the presentation
output_path = "cms_ldp_presentation.pptx"
prs.save(output_path)

output_path
